"""Turn raw CV files into plain text the rest of the pipeline can read.

What this does: reads every PDF, Word, and PowerPoint CV in a folder and
converts each one into a single block of text.

Why it exists: consultants' CVs arrive as wildly different file formats and
page layouts. Nothing downstream (trust-checking, LLM extraction, matching)
can work with a CV until it has been reduced to plain text -- and that text
has to be complete and correctly spelled, or every later step silently works
from the wrong data.

What it takes in / produces: input is a folder path containing .pdf/.docx/
.pptx files. Output is a list of dicts -- one per CV, or one per slide for a
multi-CV PowerPoint deck -- each holding which file it came from, what
format it was, the extracted text, and when it was extracted.

Assumptions and shortcuts taken:
- A file that fails to parse is skipped and logged rather than aborting the
  whole batch, since one bad CV shouldn't block everyone else's.
- Headers and footers are extracted alongside the document body, because one
  CV in this dataset hides a prompt-injection attempt in invisible
  (white-on-white) header text -- text src/trust.py can only flag if
  src/ingest.py actually captured it.
- A handful of PDF font-encoding quirks specific to this dataset (Danish
  letters and typographic ligatures rendered as unresolved font codes) are
  patched with fixes verified against the actual files, not solved in
  general -- see the comments below for the evidence behind each one.
"""

from __future__ import annotations

import logging
import re
from datetime import datetime, timezone
from pathlib import Path

import pdfminer.encodingdb as _pdfminer_encodingdb
import pdfplumber
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx"}

# Some source PDFs embed Danish letters as Type3 glyphs with non-standard
# names (e.g. "a248") and no ToUnicode CMap. pdfminer can't resolve such a
# name via the Adobe Glyph List, so it silently falls back to Adobe
# StandardEncoding for that character code -- which maps code 248 to "l with
# stroke" (ł) instead of "o with stroke" (ø), corrupting names like "Bøgh".
# These PDFs were produced with codes that match WinAnsiEncoding (cp1252),
# so for exactly the six Danish-letter codes we resolve against that table
# instead. Restricted to those codes only: unrelated custom glyphs (e.g.
# bullet icons reusing other codes) must keep falling through unresolved
# rather than being guessed at.
_DANISH_WINANSI_CODES = {197, 198, 216, 229, 230, 248}  # Å Æ Ø å æ ø
_GLYPH_NAME_CODE_RE = re.compile(r"^[A-Za-z]?(\d+)$")


def _patch_pdfminer_danish_glyph_fallback() -> None:
    """Make pdfminer resolve Danish letters (æøå) correctly in PDFs that use non-standard font glyph names."""
    winansi = _pdfminer_encodingdb.EncodingDB.encodings.get("WinAnsiEncoding", {})
    original_name2unicode = _pdfminer_encodingdb.name2unicode

    def patched_name2unicode(name: str) -> str:
        """Resolve one glyph name to a character, falling back to WinAnsi only for the known Danish-letter codes."""
        try:
            return original_name2unicode(name)
        except Exception:
            match = _GLYPH_NAME_CODE_RE.match(name.split(".")[0])
            if match:
                code = int(match.group(1))
                if code in _DANISH_WINANSI_CODES and code in winansi:
                    return winansi[code]
            raise

    _pdfminer_encodingdb.name2unicode = patched_name2unicode


_patch_pdfminer_danish_glyph_fallback()

# When a Type3 glyph has no name-based or WinAnsi-based resolution at all,
# pdfminer emits a literal "(cid:N)" placeholder rather than guessing. Across
# all 14 source PDFs there are exactly six such codes, and each one's
# identity was confirmed from its surrounding text context (e.g. "sta(cid:27)"
# next to "streamlining" -> "staff"; "2020(cid:21)Present" -> an en dash
# between dates), consistently across every file it appears in. These are
# the standard Adobe ligatures (fi/fl/ff/ffi) plus an en dash and a bullet
# glyph -- known quantities, not guesses, so it's safe to substitute them.
# Any other unresolved code is left as "(cid:N)" rather than invented.
_CID_PLACEHOLDER_SUBSTITUTIONS = {
    "21": "–",  # en dash, e.g. "2020(cid:21)Present"
    "27": "ff",  # e.g. "sta(cid:27)" -> "staff"
    "28": "fi",  # e.g. "(cid:28)rm" -> "firm"
    "29": "fl",  # e.g. "(cid:29)uent" -> "fluent"
    "30": "ffi",  # e.g. "e(cid:30)ciency" -> "efficiency"
    "136": "•",  # bullet list marker
}
_CID_PLACEHOLDER_RE = re.compile(r"\(cid:(\d+)\)")


def _resolve_known_cid_placeholders(text: str) -> str:
    """Replace pdfminer's unresolved "(cid:N)" placeholders with the real characters they stand for, where known."""

    def replace(match: re.Match[str]) -> str:
        """Look up the substitution for one matched placeholder, or leave it unchanged if unknown."""
        return _CID_PLACEHOLDER_SUBSTITUTIONS.get(match.group(1), match.group(0))

    return _CID_PLACEHOLDER_RE.sub(replace, text)


def extract_pdf_text(path: Path) -> str:
    """Pull all text out of one PDF file, page by page."""
    parts: list[str] = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                parts.append(text)
    return _resolve_known_cid_placeholders("\n".join(parts))


_HEADER_FOOTER_ATTRS = (
    "header",
    "footer",
    "first_page_header",
    "first_page_footer",
    "even_page_header",
    "even_page_footer",
)


def _header_footer_parts(part) -> list[str]:
    """Collect the visible text (paragraphs and table cells) from one document header or footer."""
    parts: list[str] = [p.text for p in part.paragraphs if p.text.strip()]
    for table in part.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    parts.append(cell.text)
    return parts


def extract_docx_text(path: Path) -> str:
    """Pull all text out of one Word file, including tables and every header/footer variant."""
    document = Document(path)
    parts: list[str] = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    parts.append(cell.text)

    # Headers/footers are separate parts that Document.paragraphs never
    # touches. A CV in this corpus hides a prompt-injection payload in
    # white-on-white header text specifically because it's invisible both
    # to a human skimming the document and to naive paragraph-only
    # extraction -- src/trust.py can't flag what src/ingest.py never sees.
    for section in document.sections:
        for attr in _HEADER_FOOTER_ATTRS:
            part = getattr(section, attr)
            if part.is_linked_to_previous:
                continue
            parts.extend(_header_footer_parts(part))

    return "\n".join(parts)


def _shape_text(shape: BaseShape) -> list[str]:
    """Collect the visible text from one PowerPoint shape, recursing into grouped shapes."""
    parts: list[str] = []
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            parts.extend(_shape_text(sub_shape))
        return parts
    if shape.has_text_frame and shape.text_frame.text.strip():
        parts.append(shape.text_frame.text)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    parts.append(cell.text)
    return parts


def extract_pptx_slides(path: Path) -> list[str]:
    """Pull all text out of one PowerPoint file, returning one text block per slide."""
    presentation = Presentation(path)
    slide_texts: list[str] = []
    for slide in presentation.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            parts.extend(_shape_text(shape))
        slide_texts.append("\n".join(parts))
    return slide_texts


def _make_record(source_file: str, fmt: str, raw_text: str) -> dict:
    """Bundle one CV's extracted text together with where it came from and when it was processed."""
    return {
        "source_file": source_file,
        "format": fmt,
        "raw_text": raw_text.strip(),
        "extracted_at": datetime.now(timezone.utc).isoformat(),
    }


def ingest_directory(directory: str | Path) -> list[dict]:
    """Convert every supported CV file in a folder into text records, skipping (and logging) any file that fails."""
    directory = Path(directory)
    if not directory.is_dir():
        raise FileNotFoundError(f"No such directory: {directory}")

    records: list[dict] = []
    for path in sorted(directory.iterdir()):
        if not path.is_file():
            continue
        suffix = path.suffix.lower()
        if suffix not in SUPPORTED_EXTENSIONS:
            continue

        try:
            if suffix == ".pdf":
                records.append(_make_record(path.name, "pdf", extract_pdf_text(path)))
            elif suffix == ".docx":
                records.append(_make_record(path.name, "docx", extract_docx_text(path)))
            elif suffix == ".pptx":
                for i, text in enumerate(extract_pptx_slides(path), start=1):
                    records.append(
                        _make_record(f"{path.name}#slide{i}", "pptx", text)
                    )
        except Exception:
            logger.exception("Failed to ingest %s", path)
            continue

    return records
